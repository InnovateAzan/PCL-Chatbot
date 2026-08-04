from __future__ import annotations

import csv
import mimetypes
import re
import shutil
import subprocess
import tempfile
import traceback
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Callable, Iterator

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pytesseract
except ImportError:
    pytesseract = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from PIL import Image, ImageEnhance, ImageFilter, ImageOps
except ImportError:
    Image = None
    ImageEnhance = None
    ImageFilter = None
    ImageOps = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


@dataclass
class LoadedPage:
    """Text extracted from one source page."""

    document_name: str
    page_number: int
    text: str
    extraction_method: str

    def to_metadata(self) -> dict[str, str | int]:
        return {
            "document_name": self.document_name,
            "page_number": self.page_number,
        }


@dataclass
class LoadedDocument:
    """Represents text extracted from one document."""

    file_name: str
    file_path: str
    file_extension: str
    mime_type: str
    text: str
    status: str
    loader_used: str
    error: str | None = None
    pages: list[LoadedPage] = field(default_factory=list)

    def to_dict(self) -> dict:
        return asdict(self)


class DocumentLoader:
    """
    Dynamic loader for policy and company documents.

    Directly supported:
    - PDF
    - Word DOCX
    - Excel XLSX and XLSM
    - PowerPoint PPTX
    - TXT, MD, JSON, XML, HTML, YAML and LOG
    - CSV
    - PNG, JPG, JPEG, BMP, TIFF and WEBP using OCR

    Unknown and legacy formats are passed to Apache Tika when enabled.
    """

    PDF_EXTENSIONS = {
        ".pdf",
    }

    WORD_EXTENSIONS = {
        ".docx",
    }

    EXCEL_EXTENSIONS = {
        ".xlsx",
        ".xlsm",
    }

    POWERPOINT_EXTENSIONS = {
        ".pptx",
    }

    CSV_EXTENSIONS = {
        ".csv",
    }

    TEXT_EXTENSIONS = {
        ".txt",
        ".md",
        ".log",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".yaml",
        ".yml",
    }

    IMAGE_EXTENSIONS = {
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
        ".tif",
        ".tiff",
        ".webp",
    }

    LEGACY_OFFICE_EXTENSIONS = {
        ".doc",
        ".xls",
        ".ppt",
        ".rtf",
        ".odt",
        ".ods",
        ".odp",
    }

    IGNORED_FILE_NAMES = {
        "thumbs.db",
        ".ds_store",
    }

    IGNORED_PREFIXES = {
        "~$",
        ".",
    }

    def __init__(
        self,
        tesseract_path: str | None = None,
        ocr_language: str = "eng",
        ocr_dpi: int = 300,
        minimum_text_length: int = 30,
        use_tika_fallback: bool = True,
    ) -> None:
        self.ocr_language = ocr_language
        self.ocr_dpi = ocr_dpi
        self.minimum_text_length = minimum_text_length
        self.use_tika_fallback = use_tika_fallback

        if tesseract_path and pytesseract:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path

        self.loaders: dict[str, Callable[[Path], str]] = {}
        self._register_default_loaders()

    def _register_default_loaders(self) -> None:
        """Register built-in loaders."""

        for extension in self.PDF_EXTENSIONS:
            self.register_loader(extension, self._load_pdf)

        for extension in self.WORD_EXTENSIONS:
            self.register_loader(extension, self._load_word)

        for extension in self.EXCEL_EXTENSIONS:
            self.register_loader(extension, self._load_excel)

        for extension in self.POWERPOINT_EXTENSIONS:
            self.register_loader(extension, self._load_powerpoint)

        for extension in self.CSV_EXTENSIONS:
            self.register_loader(extension, self._load_csv)

        for extension in self.TEXT_EXTENSIONS:
            self.register_loader(extension, self._load_text)

        for extension in self.IMAGE_EXTENSIONS:
            self.register_loader(extension, self._load_image_ocr)

    def register_loader(
        self,
        extension: str,
        loader: Callable[[Path], str],
    ) -> None:
        """Register a loader for any custom file extension."""

        normalized_extension = extension.lower().strip()

        if not normalized_extension.startswith("."):
            normalized_extension = f".{normalized_extension}"

        self.loaders[normalized_extension] = loader

    def load(self, file_path: str | Path) -> LoadedDocument:
        """Load one document and return its extracted text."""

        path = Path(file_path).resolve()

        if not path.exists():
            return self._failed_result(
                path=path,
                error=f"File does not exist: {path}",
            )

        if not path.is_file():
            return self._failed_result(
                path=path,
                error=f"Path is not a file: {path}",
            )

        if self._should_ignore(path):
            return LoadedDocument(
                file_name=path.name,
                file_path=str(path),
                file_extension=path.suffix.lower(),
                mime_type=self.detect_mime_type(path),
                text="",
                status="ignored",
                loader_used="none",
                error="Temporary or hidden file was ignored.",
            )

        extension = path.suffix.lower()
        loader = self.loaders.get(extension)
        loader_used = "unknown"

        try:
            loaded_pages: list[LoadedPage] = []

            if loader:
                loader_used = loader.__name__
                if extension in self.PDF_EXTENSIONS:
                    loaded_pages = self._load_pdf_pages(path)
                    extracted_text = "\n\n".join(
                        page.text
                        for page in loaded_pages
                    )
                else:
                    extracted_text = loader(path)

            elif self.use_tika_fallback:
                loader_used = "apache-tika"
                extracted_text = self._load_with_tika(path)

            else:
                raise ValueError(
                    f"Unsupported file format: {extension or 'unknown'}"
                )

            cleaned_text = self.clean_text(extracted_text)

            if not cleaned_text:
                raise ValueError(
                    "No readable text could be extracted from the document."
                )

            return LoadedDocument(
                file_name=path.name,
                file_path=str(path),
                file_extension=extension,
                mime_type=self.detect_mime_type(path),
                text=cleaned_text,
                status="success",
                loader_used=loader_used,
                error=None,
                pages=loaded_pages,
            )

        except Exception as error:
            print(
                f"\nDocument loading failed: {path.name}\n"
                f"Loader: {loader_used}\n"
                f"Error: {error}\n"
            )
            traceback.print_exc()

            return self._failed_result(
                path=path,
                error=str(error),
                loader_used=loader_used,
            )

    def load_folder(
        self,
        folder_path: str | Path,
        recursive: bool = True,
    ) -> list[LoadedDocument]:
        """Load all documents from a folder."""

        folder = Path(folder_path).resolve()

        if not folder.exists():
            raise FileNotFoundError(
                f"Policies folder does not exist: {folder}"
            )

        if not folder.is_dir():
            raise NotADirectoryError(
                f"Path is not a folder: {folder}"
            )

        results: list[LoadedDocument] = []

        for path in self.iter_files(
            folder_path=folder,
            recursive=recursive,
        ):
            results.append(self.load(path))

        return results

    def iter_files(
        self,
        folder_path: str | Path,
        recursive: bool = True,
    ) -> Iterator[Path]:
        """Yield valid files from a folder."""

        folder = Path(folder_path)

        paths = (
            folder.rglob("*")
            if recursive
            else folder.glob("*")
        )

        for path in sorted(paths):
            if not path.is_file():
                continue

            if self._should_ignore(path):
                continue

            yield path

    def detect_mime_type(
        self,
        file_path: str | Path,
    ) -> str:
        """Detect basic MIME type using the file name."""

        mime_type, _ = mimetypes.guess_type(str(file_path))

        return mime_type or "application/octet-stream"

    def supported_extensions(self) -> list[str]:
        """Return supported extensions."""

        extensions = set(self.loaders.keys())
        extensions.update(self.LEGACY_OFFICE_EXTENSIONS)

        return sorted(extensions)

    # ---------------------------------------------------------
    # PDF
    # ---------------------------------------------------------

    def _load_pdf(self, path: Path) -> str:
        return "\n\n".join(
            page.text
            for page in self._load_pdf_pages(path)
        )

    def _load_pdf_pages(
        self,
        path: Path,
    ) -> list[LoadedPage]:
        if fitz is None:
            raise RuntimeError(
                "PyMuPDF is not installed. "
                "Run: pip install pymupdf"
            )

        document = fitz.open(str(path))
        extracted_pages: list[LoadedPage] = []

        try:
            for page_number, page in enumerate(
                document,
                start=1,
            ):
                normal_text = page.get_text("text").strip()
                text = normal_text
                extraction_method = "text"

                if self._needs_ocr(normal_text):
                    try:
                        ocr_text = self._ocr_pdf_page(page)

                        if len(ocr_text.strip()) > len(
                            normal_text.strip()
                        ):
                            text = ocr_text
                            extraction_method = "ocr"

                    except Exception as error:
                        print(
                            f"OCR skipped for {path.name}, "
                            f"page {page_number}: {error}"
                        )

                if text.strip():
                    page_text = (
                        f"[Page {page_number} | "
                        f"{extraction_method}]\n"
                        f"{text.strip()}"
                    )

                    extracted_pages.append(
                        LoadedPage(
                            document_name=path.name,
                            page_number=page_number,
                            text=page_text,
                            extraction_method=extraction_method,
                        )
                    )

        finally:
            document.close()

        return extracted_pages

    def _ocr_pdf_page(self, page) -> str:
        """Apply OCR to one PDF page."""

        try:
            text_page = page.get_textpage_ocr(
                language=self.ocr_language,
                dpi=self.ocr_dpi,
                full=True,
            )

            return page.get_text(
                "text",
                textpage=text_page,
            ).strip()

        except Exception as pymupdf_ocr_error:
            print(
                "PyMuPDF OCR failed. Trying pytesseract: "
                f"{pymupdf_ocr_error}"
            )

        if pytesseract is None or Image is None:
            raise RuntimeError(
                "OCR dependencies are missing. Install Pillow and "
                "pytesseract, and install Tesseract OCR."
            )

        zoom = max(
            self.ocr_dpi / 72,
            2.0,
        )

        pixmap = page.get_pixmap(
            matrix=fitz.Matrix(zoom, zoom),
            alpha=False,
        )

        if pixmap.width < 3 or pixmap.height < 3:
            return ""

        image = Image.frombytes(
            "RGB",
            (pixmap.width, pixmap.height),
            pixmap.samples,
        )

        image = self._prepare_image_for_ocr(image)

        return pytesseract.image_to_string(
            image,
            lang=self.ocr_language,
            config="--oem 3 --psm 6",
            timeout=30,
        ).strip()

    # ---------------------------------------------------------
    # Image OCR
    # ---------------------------------------------------------

    def _load_image_ocr(self, path: Path) -> str:
        if pytesseract is None:
            raise RuntimeError(
                "pytesseract is not installed. "
                "Run: pip install pytesseract"
            )

        if Image is None:
            raise RuntimeError(
                "Pillow is not installed. "
                "Run: pip install pillow"
            )

        with Image.open(path) as source_image:
            image = source_image.convert("RGB")

            if image.width < 3 or image.height < 3:
                raise ValueError(
                    "Image is too small for OCR."
                )

            image = self._prepare_image_for_ocr(image)

            return pytesseract.image_to_string(
                image,
                lang=self.ocr_language,
                config="--oem 3 --psm 6",
                timeout=30,
            ).strip()

    def _prepare_image_for_ocr(self, image):
        """Improve image before OCR."""

        if (
            ImageOps is None
            or ImageEnhance is None
            or ImageFilter is None
        ):
            return image

        grayscale = ImageOps.grayscale(image)
        grayscale = ImageOps.autocontrast(grayscale)

        sharpness = ImageEnhance.Sharpness(grayscale)
        grayscale = sharpness.enhance(1.5)

        contrast = ImageEnhance.Contrast(grayscale)
        grayscale = contrast.enhance(1.3)

        grayscale = grayscale.filter(
            ImageFilter.SHARPEN
        )

        return grayscale

    # ---------------------------------------------------------
    # Word
    # ---------------------------------------------------------

    def _load_word(self, path: Path) -> str:
        if Document is None:
            raise RuntimeError(
                "python-docx is not installed. "
                "Run: pip install python-docx"
            )

        rendered_text = self._load_word_as_paginated_pdf(path)

        if rendered_text:
            return rendered_text

        document = Document(str(path))
        content: list[str] = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()

            if text:
                content.append(text)

        for table_number, table in enumerate(
            document.tables,
            start=1,
        ):
            content.append(
                f"[Table {table_number}]"
            )

            for row in table.rows:
                values = [
                    self._normalize_cell_value(cell.text)
                    for cell in row.cells
                ]

                if any(values):
                    content.append(
                        " | ".join(values)
                    )

        for section_number, section in enumerate(
            document.sections,
            start=1,
        ):
            header_lines = [
                paragraph.text.strip()
                for paragraph in section.header.paragraphs
                if paragraph.text.strip()
            ]

            footer_lines = [
                paragraph.text.strip()
                for paragraph in section.footer.paragraphs
                if paragraph.text.strip()
            ]

            if header_lines:
                content.append(
                    f"[Header {section_number}]\n"
                    + "\n".join(header_lines)
                )

            if footer_lines:
                content.append(
                    f"[Footer {section_number}]\n"
                    + "\n".join(footer_lines)
                )

        return "\n".join(content)

    def _load_word_as_paginated_pdf(
        self,
        path: Path,
    ) -> str | None:
        temp_dir = Path(
            tempfile.mkdtemp(prefix="pcl-word-pages-")
        )
        pdf_path = temp_dir / f"{path.stem}.pdf"

        try:
            converted = (
                self._convert_word_to_pdf_with_word(path, pdf_path)
                or self._convert_word_to_pdf_with_libreoffice(
                    path,
                    temp_dir,
                )
            )

            if not converted or not pdf_path.exists():
                return None

            paginated_text = self._load_pdf(pdf_path).strip()

            return paginated_text or None

        except Exception:
            return None

        finally:
            shutil.rmtree(
                temp_dir,
                ignore_errors=True,
            )

    def _convert_word_to_pdf_with_word(
        self,
        source_path: Path,
        pdf_path: Path,
    ) -> bool:
        powershell = shutil.which("powershell")

        if not powershell:
            return False

        source_literal = str(source_path).replace(
            "'",
            "''",
        )
        pdf_literal = str(pdf_path).replace(
            "'",
            "''",
        )

        script = "\n".join(
            [
                "$ErrorActionPreference = 'Stop'",
                "$word = $null",
                "$document = $null",
                "try {",
                "  $word = New-Object -ComObject Word.Application",
                "  $word.Visible = $false",
                "  $document = $word.Documents.Open("
                f"'{source_literal}', $false, $true)",
                "  $document.SaveAs("
                f"'{pdf_literal}', 17)",
                "} finally {",
                "  if ($document) { $document.Close($false) }",
                "  if ($word) { $word.Quit() }",
                "}",
            ]
        )

        result = subprocess.run(
            [
                powershell,
                "-NoProfile",
                "-NonInteractive",
                "-Command",
                script,
            ],
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
        )

        return (
            result.returncode == 0
            and pdf_path.exists()
        )

    def _convert_word_to_pdf_with_libreoffice(
        self,
        source_path: Path,
        output_dir: Path,
    ) -> bool:
        soffice = shutil.which("soffice")

        if not soffice:
            return False

        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(output_dir),
                str(source_path),
            ],
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
        )

        expected_pdf = output_dir / f"{source_path.stem}.pdf"

        return (
            result.returncode == 0
            and expected_pdf.exists()
        )

    # ---------------------------------------------------------
    # Excel
    # ---------------------------------------------------------

    def _load_excel(self, path: Path) -> str:
        if openpyxl is None:
            raise RuntimeError(
                "openpyxl is not installed. "
                "Run: pip install openpyxl"
            )

        workbook = openpyxl.load_workbook(
            filename=str(path),
            read_only=True,
            data_only=True,
        )

        content: list[str] = []

        try:
            for worksheet in workbook.worksheets:
                content.append(
                    f"[Worksheet: {worksheet.title}]"
                )

                for row_number, row in enumerate(
                    worksheet.iter_rows(
                        values_only=True
                    ),
                    start=1,
                ):
                    values = [
                        self._normalize_cell_value(value)
                        for value in row
                    ]

                    while values and not values[-1]:
                        values.pop()

                    if not any(values):
                        continue

                    content.append(
                        f"[Row {row_number}] "
                        + " | ".join(values)
                    )

        finally:
            workbook.close()

        return "\n".join(content)

    # ---------------------------------------------------------
    # PowerPoint
    # ---------------------------------------------------------

    def _load_powerpoint(self, path: Path) -> str:
        if Presentation is None:
            raise RuntimeError(
                "python-pptx is not installed. "
                "Run: pip install python-pptx"
            )

        presentation = Presentation(str(path))
        content: list[str] = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            content.append(
                f"[Slide {slide_number}]"
            )

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()

                    if text:
                        content.append(text)

                if getattr(
                    shape,
                    "has_table",
                    False,
                ):
                    for row in shape.table.rows:
                        values = [
                            self._normalize_cell_value(
                                cell.text
                            )
                            for cell in row.cells
                        ]

                        if any(values):
                            content.append(
                                " | ".join(values)
                            )

            try:
                if slide.has_notes_slide:
                    notes_text = (
                        slide.notes_slide
                        .notes_text_frame
                        .text
                        .strip()
                    )

                    if notes_text:
                        content.append(
                            f"[Slide {slide_number} Notes]\n"
                            f"{notes_text}"
                        )

            except Exception:
                pass

        return "\n".join(content)

    # ---------------------------------------------------------
    # Text and CSV
    # ---------------------------------------------------------

    def _load_text(self, path: Path) -> str:
        encodings = (
            "utf-8",
            "utf-8-sig",
            "cp1252",
            "latin-1",
        )

        for encoding in encodings:
            try:
                return path.read_text(
                    encoding=encoding,
                    errors="strict",
                )

            except UnicodeDecodeError:
                continue

        return path.read_text(
            encoding="utf-8",
            errors="ignore",
        )

    def _load_csv(self, path: Path) -> str:
        content: list[str] = []
        encoding = self._detect_text_encoding(path)

        with path.open(
            mode="r",
            encoding=encoding,
            errors="ignore",
            newline="",
        ) as csv_file:
            sample = csv_file.read(4096)
            csv_file.seek(0)

            try:
                dialect = csv.Sniffer().sniff(sample)

            except csv.Error:
                dialect = csv.excel

            reader = csv.reader(
                csv_file,
                dialect,
            )

            for row_number, row in enumerate(
                reader,
                start=1,
            ):
                values = [
                    self._normalize_cell_value(value)
                    for value in row
                ]

                if any(values):
                    content.append(
                        f"[Row {row_number}] "
                        + " | ".join(values)
                    )

        return "\n".join(content)

    # ---------------------------------------------------------
    # Apache Tika
    # ---------------------------------------------------------

    def _load_with_tika(self, path: Path) -> str:
        """Use Apache Tika for unknown and legacy formats."""

        try:
            from tika import parser

        except ImportError as error:
            raise ValueError(
                f"Unsupported file format: "
                f"{path.suffix or 'unknown'}. "
                "Install Tika using: pip install tika"
            ) from error

        parsed = parser.from_file(str(path))

        if not parsed:
            return ""

        content = parsed.get("content") or ""

        return str(content).strip()

    # ---------------------------------------------------------
    # Utilities
    # ---------------------------------------------------------

    def _needs_ocr(self, text: str) -> bool:
        """
        Use OCR only when normal PDF extraction does not contain
        meaningful readable text.
        """

        if not text:
            return True

        cleaned = self.clean_text(text)

        if not cleaned:
            return True

        alphabetic_count = sum(
            character.isalpha()
            for character in cleaned
        )

        word_count = len(cleaned.split())

        if alphabetic_count >= 10:
            return False

        if word_count >= 3:
            return False

        return True

    def _should_ignore(self, path: Path) -> bool:
        file_name = path.name.lower()

        if file_name in self.IGNORED_FILE_NAMES:
            return True

        return any(
            path.name.startswith(prefix)
            for prefix in self.IGNORED_PREFIXES
        )

    def _failed_result(
        self,
        path: Path,
        error: str,
        loader_used: str = "none",
    ) -> LoadedDocument:
        return LoadedDocument(
            file_name=path.name,
            file_path=str(path),
            file_extension=path.suffix.lower(),
            mime_type=self.detect_mime_type(path),
            text="",
            status="failed",
            loader_used=loader_used,
            error=error,
        )

    @staticmethod
    def _normalize_cell_value(
        value: object,
    ) -> str:
        if value is None:
            return ""

        text = str(value).strip()

        return re.sub(
            r"\s+",
            " ",
            text,
        )

    @staticmethod
    def clean_text(text: str) -> str:
        """Clean extracted text."""

        if not text:
            return ""

        text = text.replace(
            "\x00",
            " ",
        )
        text = text.replace(
            "\r\n",
            "\n",
        )
        text = text.replace(
            "\r",
            "\n",
        )

        lines: list[str] = []

        for raw_line in text.splitlines():
            line = re.sub(
                r"[ \t]+",
                " ",
                raw_line,
            ).strip()

            lines.append(line)

        cleaned_lines: list[str] = []
        previous_was_blank = False

        for line in lines:
            is_blank = not line

            if is_blank and previous_was_blank:
                continue

            cleaned_lines.append(line)
            previous_was_blank = is_blank

        return "\n".join(
            cleaned_lines
        ).strip()

    @staticmethod
    def _detect_text_encoding(
        path: Path,
    ) -> str:
        encodings = (
            "utf-8-sig",
            "utf-8",
            "cp1252",
            "latin-1",
        )

        for encoding in encodings:
            try:
                with path.open(
                    "r",
                    encoding=encoding,
                    errors="strict",
                ) as file:
                    file.read(4096)

                return encoding

            except UnicodeDecodeError:
                continue

        return "utf-8"


def load_policy_folder(
    folder_path: str | Path,
    tesseract_path: str | None = None,
    ocr_language: str = "eng",
) -> list[dict[str, str | None]]:
    """Load all documents from the policies folder."""

    loader = DocumentLoader(
        tesseract_path=tesseract_path,
        ocr_language=ocr_language,
    )

    documents = loader.load_folder(
        folder_path=folder_path,
        recursive=True,
    )

    return [
        document.to_dict()
        for document in documents
    ]


# Important:
# Existing policies.py imports PolicyDocumentLoader.
# This alias must remain outside the __main__ block.
PolicyDocumentLoader = DocumentLoader


if __name__ == "__main__":
    project_root = Path(
        __file__
    ).resolve().parents[3]

    policies_folder = (
        project_root / "policies"
    )

    loader = DocumentLoader(
        tesseract_path=(
            r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        ),
        ocr_language="eng",
    )

    print(
        f"Reading documents from: "
        f"{policies_folder}"
    )

    print(
        "Supported formats: "
        f"{', '.join(loader.supported_extensions())}"
    )

    try:
        loaded_documents = loader.load_folder(
            folder_path=policies_folder,
            recursive=True,
        )

        successful_count = 0
        failed_count = 0
        ignored_count = 0

        for document in loaded_documents:
            print(
                f"{document.status.upper():8} | "
                f"{document.file_name} | "
                f"{document.loader_used}"
            )

            if document.error:
                print(
                    f"Error: {document.error}"
                )

            if document.status == "success":
                successful_count += 1

            elif document.status == "failed":
                failed_count += 1

            elif document.status == "ignored":
                ignored_count += 1

        print(
            "\nDocument loading completed."
        )
        print(
            f"Successful: {successful_count}"
        )
        print(
            f"Failed: {failed_count}"
        )
        print(
            f"Ignored: {ignored_count}"
        )

    except Exception as error:
        print(
            f"Folder loading error: {error}"
        )
        traceback.print_exc()
